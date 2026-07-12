import hashlib
import json
import logging
import mimetypes
from collections import OrderedDict
from pathlib import Path
from typing import Any

from langchain_core.documents import Document
from langchain_ollama import OllamaEmbeddings

from app.core.config import settings

logger = logging.getLogger(__name__)


TEXT_EXTENSIONS = {
    ".txt",
    ".md",
    ".markdown",
    ".py",
    ".ts",
    ".tsx",
    ".js",
    ".jsx",
    ".json",
    ".xml",
    ".csv",
    ".html",
    ".css",
    ".scss",
    ".sql",
    ".yml",
    ".yaml",
    ".toml",
}
IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".webp", ".gif", ".bmp", ".tiff"}

POWER_BI_EXTENSIONS = {".pbix", ".pbit"}

class BoundedCache:
    def __init__(self, max_size: int) -> None:
        self.max_size = max(1, max_size)
        self._items: OrderedDict[str, Any] = OrderedDict()

    def get(self, key: str) -> Any | None:
        value = self._items.get(key)
        if value is not None:
            self._items.move_to_end(key)
        return value

    def set(self, key: str, value: Any) -> None:
        self._items[key] = value
        self._items.move_to_end(key)
        while len(self._items) > self.max_size:
            self._items.popitem(last=False)


class MultimodalExtractor:
    def extract(self, path: Path, content_type: str) -> tuple[str, list[str], dict[str, Any]]:
        extension = path.suffix.lower()
        metadata: dict[str, Any] = {"extension": extension}

        if extension == ".pdf":
            return self._extract_pdf(path, metadata)
        if extension == ".docx":
            return self._extract_docx(path, metadata)
        if extension in {".xlsx", ".xlsm", ".xls"}:
            return self._extract_spreadsheet(path, metadata)
        if extension == ".pptx":
            return self._extract_presentation(path, metadata)
        if extension in POWER_BI_EXTENSIONS:
            return self._extract_power_bi(path, metadata)
        if extension in IMAGE_EXTENSIONS or content_type.startswith("image/"):
            return self._extract_image(path, metadata)
        if extension in TEXT_EXTENSIONS or content_type.startswith("text/"):
            return self._extract_text(path, metadata)

        try:
            return self._extract_text(path, metadata)
        except UnicodeDecodeError:
            metadata["warning"] = "Formato binário sem extrator textual configurado."
            return "", ["binary"], metadata

    def _extract_text(self, path: Path, metadata: dict[str, Any]) -> tuple[str, list[str], dict[str, Any]]:
        text = path.read_text(encoding="utf-8")
        metadata["lines"] = text.count("\n") + 1
        return text, ["text"], metadata

    def _extract_pdf(self, path: Path, metadata: dict[str, Any]) -> tuple[str, list[str], dict[str, Any]]:
        try:
            from pypdf import PdfReader
        except ImportError:
            metadata["warning"] = "pypdf não está instalado."
            return "", ["pdf"], metadata

        reader = PdfReader(str(path))
        pages = []
        for index, page in enumerate(reader.pages):
            text = page.extract_text() or ""
            pages.append(f"[Página {index + 1}]\n{text}")
        metadata["pages"] = len(reader.pages)
        return "\n\n".join(pages), ["pdf", "text"], metadata

    def _extract_docx(self, path: Path, metadata: dict[str, Any]) -> tuple[str, list[str], dict[str, Any]]:
        try:
            from docx import Document as DocxDocument
        except ImportError:
            metadata["warning"] = "python-docx não está instalado."
            return "", ["docx"], metadata

        doc = DocxDocument(str(path))
        paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
        table_rows = []
        for table in doc.tables:
            for row in table.rows:
                table_rows.append(" | ".join(cell.text.strip() for cell in row.cells))
        metadata["paragraphs"] = len(paragraphs)
        metadata["tables"] = len(doc.tables)
        return "\n".join(paragraphs + table_rows), ["docx", "text", "table"], metadata

    def _extract_spreadsheet(self, path: Path, metadata: dict[str, Any]) -> tuple[str, list[str], dict[str, Any]]:
        try:
            from openpyxl import load_workbook
        except ImportError:
            metadata["warning"] = "openpyxl não está instalado."
            return "", ["spreadsheet"], metadata

        workbook = load_workbook(str(path), data_only=True, read_only=True)
        rows: list[str] = []
        for sheet in workbook.worksheets:
            rows.append(f"[Planilha: {sheet.title}]")
            for row in sheet.iter_rows(values_only=True):
                values = ["" if value is None else str(value) for value in row]
                if any(value.strip() for value in values):
                    rows.append(" | ".join(values))
        metadata["sheets"] = [sheet.title for sheet in workbook.worksheets]
        return "\n".join(rows), ["spreadsheet", "table", "text"], metadata

    def _extract_presentation(self, path: Path, metadata: dict[str, Any]) -> tuple[str, list[str], dict[str, Any]]:
        try:
            from pptx import Presentation
        except ImportError:
            metadata["warning"] = "python-pptx não está instalado."
            return "", ["presentation"], metadata

        presentation = Presentation(str(path))
        slides: list[str] = []
        for index, slide in enumerate(presentation.slides):
            texts = []
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    texts.append(shape.text.strip())
            slides.append(f"[Slide {index + 1}]\n" + "\n".join(texts))
        metadata["slides"] = len(presentation.slides)
        return "\n\n".join(slides), ["presentation", "diagram", "text"], metadata

    def _extract_image(self, path: Path, metadata: dict[str, Any]) -> tuple[str, list[str], dict[str, Any]]:
        try:
            from PIL import Image
        except ImportError:
            metadata["warning"] = "Pillow não está instalado."
            return "", ["image"], metadata

        with Image.open(path) as image:
            metadata["width"] = image.width
            metadata["height"] = image.height
            metadata["mode"] = image.mode

            ocr_text = ""
            try:
                import pytesseract

                ocr_text = pytesseract.image_to_string(image, lang="por+eng").strip()
            except Exception as exc:  # OCR is optional and depends on local Tesseract binaries.
                metadata["ocrWarning"] = str(exc)

        visual_summary = (
            f"Imagem/captura/diagrama enviado: {path.name}. "
            f"Dimensões: {metadata.get('width')}x{metadata.get('height')}. "
        )
        if ocr_text:
            visual_summary += f"Texto OCR extraído:\n{ocr_text}"
            return visual_summary, ["image", "ocr", "diagram"], metadata
        return visual_summary, ["image", "diagram"], metadata
    def _extract_power_bi(self, path: Path, metadata: dict[str, Any]) -> tuple[str, list[str], dict[str, Any]]:
        # PBIX/PBIT are binary Power BI containers. Keep metadata rather than
        # attempting a UTF-8 decode; a dedicated connector can enrich it later.
        metadata["sizeBytes"] = path.stat().st_size
        metadata["warning"] = "PBIX/PBIT armazenado; a extração do modelo requer conector Power BI."
        return (
            f"Arquivo Power BI: {path.name}. Tipo: {path.suffix[1:].upper()}.",
            ["power-bi", "binary"],
            metadata,
        )


class RAGEngine:
    def __init__(self) -> None:
        self.persist_dir = Path(settings.chroma_persist_dir)
        self.upload_dir = Path(settings.rag_upload_dir)
        self.persist_dir.mkdir(parents=True, exist_ok=True)
        self.upload_dir.mkdir(parents=True, exist_ok=True)
        self.manifest_path = self.persist_dir / "manifest.json"
        self.extractor = MultimodalExtractor()
        self.response_cache = BoundedCache(settings.rag_cache_size)
        self._manifest = self._load_manifest()
        self._vectorstore = None
        self._splitter = None

    @property
    def vectorstore(self):
        if self._vectorstore is None:
            from langchain_chroma import Chroma

            embeddings = OllamaEmbeddings(
                model=settings.ollama_embedding_model,
                base_url=settings.ollama_base_url,
            )
            self._vectorstore = Chroma(
                collection_name=settings.rag_collection_name,
                embedding_function=embeddings,
                persist_directory=str(self.persist_dir),
            )
        return self._vectorstore

    @property
    def splitter(self):
        if self._splitter is None:
            from langchain_text_splitters import RecursiveCharacterTextSplitter

            self._splitter = RecursiveCharacterTextSplitter(
                chunk_size=settings.rag_chunk_size,
                chunk_overlap=settings.rag_chunk_overlap,
                separators=["\n\n", "\n", ". ", " ", ""],
            )
        return self._splitter

    def ingest_bytes(self, filename: str, content: bytes, content_type: str | None = None) -> dict[str, Any]:
        content_hash = hashlib.sha256(content).hexdigest()
        known = self._manifest.get(content_hash)
        if known:
            return {**known, "deduplicated": True}

        safe_name = Path(filename).name or f"{content_hash}.txt"
        content_type = content_type or mimetypes.guess_type(safe_name)[0] or "application/octet-stream"
        stored_path = self.upload_dir / f"{content_hash[:16]}-{safe_name}"
        stored_path.write_bytes(content)

        extracted_text, modalities, extraction_metadata = self.extractor.extract(stored_path, content_type)
        if not extracted_text.strip():
            extracted_text = f"Arquivo {safe_name} sem texto extraível automaticamente."

        document_id = content_hash[:16]
        base_metadata = {
            "document_id": document_id,
            "source": safe_name,
            "stored_path": str(stored_path),
            "content_type": content_type,
            "content_hash": content_hash,
            "source_priority": 1,
            "source_scope": "uploaded",
            "modalities": ",".join(modalities),
            **extraction_metadata,
        }
        chunks = self.splitter.split_text(extracted_text)
        documents = [
            Document(
                page_content=chunk,
                metadata={**base_metadata, "chunk": index},
            )
            for index, chunk in enumerate(chunks)
        ]
        ids = [f"{document_id}:{index}" for index in range(len(documents))]
        self.vectorstore.add_documents(documents, ids=ids)

        record = {
            "documentId": document_id,
            "filename": safe_name,
            "contentType": content_type,
            "chunks": len(documents),
            "deduplicated": False,
            "modalities": modalities,
        }
        self._manifest[content_hash] = record
        self._save_manifest()
        return record

    def retrieve(self, query: str, document_ids: list[str] | None = None, top_k: int | None = None) -> list[Document]:
        top_k = top_k or settings.rag_top_k
        docs: list[Document] = []

        if document_ids:
            for document_id in document_ids:
                docs.extend(
                    self._similarity_search(
                        query,
                        top_k,
                        {"document_id": document_id, "source_scope": "uploaded"},
                    )
                )

        docs.extend(self._similarity_search(query, top_k, None))
        unique: dict[str, Document] = {}
        for doc in docs:
            key = f"{doc.metadata.get('document_id')}:{doc.metadata.get('chunk')}:{doc.page_content[:80]}"
            unique[key] = doc

        return sorted(
            unique.values(),
            key=lambda doc: int(doc.metadata.get("source_priority", 99)),
        )[:top_k]

    def list_documents(self) -> list[dict[str, Any]]:
        """List indexed documents without exposing internal storage paths."""
        records = list(self._manifest.values())
        return sorted(records, key=lambda item: item.get("filename", "").lower())

    def cache_key(self, message: str, context: str, document_ids: list[str]) -> str:
        normalized = json.dumps(
            {
                "message": " ".join(message.lower().split()),
                "context": " ".join(context.lower().split())[:1000],
                "documentIds": sorted(document_ids),
            },
            sort_keys=True,
        )
        return hashlib.sha256(normalized.encode("utf-8")).hexdigest()

    def _similarity_search(self, query: str, top_k: int, metadata_filter: dict[str, Any] | None) -> list[Document]:
        try:
            if metadata_filter:
                return self.vectorstore.similarity_search(query, k=top_k, filter=metadata_filter)
            return self.vectorstore.similarity_search(query, k=top_k)
        except Exception as exc:
            logger.warning("Chroma retrieval failed: %s", exc)
            return []

    def _load_manifest(self) -> dict[str, Any]:
        if not self.manifest_path.exists():
            return {}
        try:
            return json.loads(self.manifest_path.read_text(encoding="utf-8"))
        except json.JSONDecodeError:
            logger.warning("Ignoring invalid RAG manifest at %s", self.manifest_path)
            return {}

    def _save_manifest(self) -> None:
        self.manifest_path.write_text(
            json.dumps(self._manifest, ensure_ascii=False, indent=2),
            encoding="utf-8",
        )
